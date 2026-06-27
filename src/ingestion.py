import os
import json
import csv
from pathlib import Path
from typing import List, Dict, Any


def _load_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def _load_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text:
            pages.append(f"[Página {i}]\n{text}")
    return "\n\n".join(pages)


def _load_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _load_excel(file_path: str) -> str:
    import pandas as pd
    dfs = []
    xl = pd.ExcelFile(file_path)
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        dfs.append(f"[Hoja: {sheet}]\n{df.to_markdown(index=False)}")
    return "\n\n".join(dfs)


def _load_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"[Diapositiva {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _load_csv(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        rows = []
        for row in reader:
            row_text = " | ".join(f"{k}: {v}" for k, v in row.items())
            rows.append(row_text)
    return "\n".join(rows)


def _load_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False, indent=2)


def _load_html(file_path: str) -> str:
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    return soup.get_text(separator="\n", strip=True)


_LOADERS = {
    ".md": _load_text,
    ".txt": _load_text,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".doc": _load_docx,
    ".xlsx": _load_excel,
    ".xls": _load_excel,
    ".pptx": _load_pptx,
    ".ppt": _load_pptx,
    ".csv": _load_csv,
    ".json": _load_json,
    ".html": _load_html,
    ".htm": _load_html,
}

SUPPORTED_EXTENSIONS = set(_LOADERS.keys())


def _infer_category(filename: str) -> str:
    name = filename.lower()
    if "privacidad" in name or "datos" in name:
        return "Privacidad y Datos"
    if "faq" in name or "preguntas" in name or "frecuentes" in name:
        return "FAQ"
    if "cancelacion" in name or "reagendamiento" in name:
        return "Políticas"
    if "convenio" in name or "cobertura" in name or "eps" in name:
        return "Convenios y Coberturas"
    if "instruccion" in name or "consulta" in name or "pre" in name or "post" in name:
        return "Instrucciones Médicas"
    return "General"


def load_document(file_path: str) -> Dict[str, Any]:
    path = Path(file_path)
    ext = path.suffix.lower()
    if ext not in _LOADERS:
        raise ValueError(f"Formato no soportado: {ext}")
    content = _LOADERS[ext](file_path)
    return {
        "content": content,
        "metadata": {
            "source": path.name,
            "file_path": str(file_path),
            "file_type": ext.lstrip("."),
            "category": _infer_category(path.name),
        },
    }


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 120) -> List[str]:
    chunks: List[str] = []
    start = 0
    length = len(text)

    while start < length:
        end = min(start + chunk_size, length)

        if end < length:
            para = text.rfind("\n\n", start, end)
            if para > start + chunk_size // 2:
                end = para + 2
            else:
                sent = text.rfind(". ", start, end)
                if sent > start + chunk_size // 2:
                    end = sent + 2

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= length:
            break
        start = end - overlap

    return chunks


def load_and_chunk_documents(docs_dir: str) -> List[Dict[str, Any]]:
    docs_path = Path(docs_dir)
    all_chunks: List[Dict[str, Any]] = []

    for file_path in sorted(docs_path.rglob("*")):
        if not file_path.is_file():
            continue
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue

        try:
            print(f"  Cargando: {file_path.name}")
            doc = load_document(str(file_path))
            chunks = chunk_text(doc["content"])

            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    "content": chunk,
                    "metadata": {
                        **doc["metadata"],
                        "chunk_index": i,
                        "total_chunks": len(chunks),
                    },
                })
            print(f"    -> {len(chunks)} fragmentos")
        except Exception as exc:
            import traceback
            print(f"    [ADVERTENCIA] No se pudo cargar {file_path.name}: {exc!r}")
            traceback.print_exc()

    return all_chunks
